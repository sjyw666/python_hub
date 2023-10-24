#研发者：时间遗忘
#开发时间：2023/10/9 10:09
import os
import docx
from PyPDF2 import PdfFileReader
from pptx import Presentation

# 指定要查找的目录和字符串
directory = '.'  # 当前路径
search_string = '指定的字符串'

# 用于存储包含指定字符串的文件列表
matching_files = []

# 遍历目录中的文件
for root, dirs, files in os.walk(directory):
    for file in files:
        if file.endswith('.docx'):
            # 处理Word文档
            docx_path = os.path.join(root, file)
            doc = docx.Document(docx_path)
            for paragraph in doc.paragraphs:
                if search_string in paragraph.text:
                    matching_files.append(docx_path)
                    break
        elif file.endswith('.pdf'):
            # 处理PDF文件
            pdf_path = os.path.join(root, file)
            pdf = PdfFileReader(open(pdf_path, 'rb'))
            for page_num in range(pdf.getNumPages()):
                page = pdf.getPage(page_num)
                page_text = page.extractText()
                if search_string in page_text:
                    matching_files.append(pdf_path)
                    break
        elif file.endswith('.pptx'):
            # 处理PPTX文件
            pptx_path = os.path.join(root, file)
            presentation = Presentation(pptx_path)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if search_string in shape.text:
                            matching_files.append(pptx_path)
                            break

# 输出包含指定字符串的文件列表和数量
print("包含指定字符串的文件列表:")
for file in matching_files:
    print(file)
print(f"总共找到了 {len(matching_files)} 个文件包含指定字符串。")
