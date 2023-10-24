#研发者：时间遗忘
#开发时间：2023/9/3 16:21
import os
import time
import PyPDF2
from docx import Document
import pptx
import rarfile


def search_files(directory, target_string):
    found_files = []
    for root, dirs, files in os.walk(directory):
        for file in files:
            file_path = os.path.join(root, file)
            file_ext = os.path.splitext(file)[1].lower()

            if file_ext == ".pdf":
                with open(file_path, "rb") as f:
                    pdf_reader = PyPDF2.PdfReader(f)
                    for page in pdf_reader.pages:
                        text = page.extract_text()
                        if target_string in text:
                            found_files.append(file_path)
                            break

            elif file_ext == ".docx":
                doc = Document(file_path)
                for para in doc.paragraphs:
                    if target_string in para.text:
                        found_files.append(file_path)
                        break

            elif file_ext == ".pptx":
                ppt = pptx.Presentation(file_path)
                for slide in ppt.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            if target_string in shape.text:
                                found_files.append(file_path)
                                break

            elif file_ext == ".rar":
                with rarfile.RarFile(file_path) as rf:
                    for f in rf.infolist():
                        if not f.isdir():
                            with rf.open(f) as subfile:
                                content = subfile.read().decode("utf-8", errors="ignore")
                                if target_string in content:
                                    found_files.append(file_path)
                                    break

    return found_files

if __name__ == "__main__":
    start_time = time.time()

    directory = input('搜索的文件夹路径：')  # 指定要搜索的文件夹路径
    target_string = input('要查找的字符串：')  # 指定要查找的字符串

    result_files = search_files(directory, target_string)

    for file in result_files:
        print(file)

    end_time = time.time()
    execution_time = end_time - start_time
    print("程序运行时间：", execution_time, "秒")
